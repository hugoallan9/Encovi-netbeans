# -*- coding: utf-8 -*-
# To change this license header, choose License Headers in Project Properties.
# To change this template file, choose Tools | Templates
# and open the template in the editor.

__author__ = "hugo"
__date__ = "$23/11/2016 04:42:04 PM$"


import glob, os
import shutil as sh
import subprocess
import time
import csv
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
import unicodedata

class Document:

    def __init__(self, tituloDoc, lugar_geografico, ruta_salida, anio = 2014):
        self.titulo_documento = tituloDoc
        self.lugar_geografico = lugar_geografico
        self.ruta_salida = ruta_salida
        self.ruta_salida = unicode(self.ruta_salida, 'utf8')
        self.ruta_salida = unicodedata.normalize("NFKD", self.ruta_salida)
        self.ruta_salida = self.ruta_salida.encode("ascii", "ignore")
        print self.ruta_salida
        self.anio_evento = anio
        self.no_capitulos = []
        self.capitulos = []
        self.titulo_seccion = []
        self.descripcion = []
        self.titulo_grafica = []
        self.desagregacion_grafica = []
        self.grafica = []
        self.fuente = []
        self.tipo_descriptor = []
        self.incluir_presentacion = []
        self.tabla = ''
        self.datos_depto = []

    def crear_directorio(self):

        try:
            os.makedirs(self.ruta_salida)
        except OSError:
            print "El directorio ya existe"
        try:
            os.makedirs(os.path.join(self.ruta_salida,'graficasPresentacion'))
        except OSError:
            print "El directorio ya existe"

    def copiar_utilidades(self):
        errors = []
        src = os.path.join(os.getcwd(),"Utilidades")
        names = os.listdir( src )
        for name in names:
            srcname = os.path.join(src , name )
            dstname = os.path.join(self.ruta_salida, name)
            try:
                sh.copy(srcname, dstname)
            except (IOError, os.error) as why:
                errors.append((srcname, dstname, str(why)))
            except Error as err:
                errors.extend(err.args[0])



    def crear_documento(self):
        self.ruta_compilacion = self.ruta_salida.strip().replace(" ", "\\ ")
        self.documento = open( os.path.join( self.ruta_salida,
        self.titulo_documento + '.tex' ), 'w+' )
        self.documento.write('%Creado de manera automática en ' + time.strftime("%x")
        + " a las " + time.strftime("%X") + '\n')
        self.documento.write('\\input{Carta3.tex} \n')
        self.documento.write('\\renewcommand{\partes}{} \n')
        self.documento.write('\\renewcommand{\\titulodoc}{ ' +
        self.titulo_documento + '}\n')
        self.documento.write( '\\newcommand{\\ra}[1]{\\renewcommand{\\arraystretch}{#1}} \n' )
        self.documento.write( '\\definecolor{color1}{rgb}{0,0,0.8} \n' +
        '\\definecolor{color2}{rgb}{0.3,0.5,1} \n' )
        self.documento.write('\\begin{document} \n' )
        self.documento.write('\\includepdf{caratula.pdf}')
        self.documento.write('\input{participantes.tex}\n')
        self.documento.write('\\tableofcontents')
        self.documento.write('\pagestyle{estandar}\n')
        self.documento.write('\setcounter{page}{0}\n')

    def crear_caratula(self):
        self.ruta_compilacion = self.ruta_salida.strip().replace(" ", "\\ ")
        documento = open( os.path.join( self.ruta_salida,
        'depto.tex' ), 'w+' )
        documento.write('Departamento de ' + self.lugar_geografico)
        documento.close()
        cadena_compilacion = "cd "+ self.ruta_compilacion + " && xelatex " + "caratula.tex"
        print cadena_compilacion
        print subprocess.Popen(cadena_compilacion, shell=True, stdout=subprocess.PIPE).stdout.read()

    def limpiar_directorio(self):
        for f in os.listdir(self.ruta_salida):
            if f == self.titulo_documento + '.pdf' or f == self.titulo_documento + '-Presentacion.pdf':
                pass
            else:
                print f
                if os.path.isfile( os.path.join(self.ruta_salida,f)):
                    os.remove(os.path.join(self.ruta_salida,f) )
                else:
                    sh.rmtree(os.path.join(self.ruta_salida,f))


    def crear_presentacion(self):
        self.presentacion = open( os.path.join( self.ruta_salida, self.titulo_documento + '-Presentacion.tex' ), 'w+' )
        self.presentacion.write('%Creado de manera automática en ' + time.strftime("%x")
        + " a las " + time.strftime("%X") + '\n')
        self.presentacion.write('\\input{Presentacion.tex} \n')
        self.presentacion.write('\\newcommand{\\ra}[1]{\\renewcommand{\\arraystretch}{#1}}')
        self.presentacion.write('\\begin{document} \n' )
        self.presentacion.write('\\primeradiapositiva{ Resultados ENCOVI 2014}{'\
        +'Participación ciudadana y medios de comunicación}{}{' + self.lugar_geografico+', noviembre 2016}')
        self.presentacion.write('\\diaposimple{Objetivo General de ENCOVI}{ Conocer y evaluar las condiciones de vida de la población y determinar los niveles de pobreza existentes en Guatemala. }\n\n')
        self.presentacion.write('\\diapolist{Objetivos Específicos de ENCOVI}{%\ \n \\item Contar con información confiable y oportuna que permita identificar las condiciones de vida de los distintos grupos sociales del país, especialmente en la estructura de los ingresos y gastos del hogar, que faciliten la elaboración y evaluación de planes, políticas y estrategias de desarrollo. \n'\
        +  '\\item Obtener estimaciones de la tasa de pobreza y pobreza extrema para cada uno de los dominios de estudio de esta encuesta. \n'\
        +'\\item Generar información socio-demográfica y económica que permita aproximarse a los niveles de bienestar de las familias y explicar sus hábitos de consumo y la manera en la que se forma su ingreso. \n'\
        + '\\item Monitorear los avances e impactos de los programas y acciones sociales. }\n\n')
        self.presentacion.write('\\diaposimple{Muestra Encovi por Departamentos}{\\normalsize\n'\
        +'\\begin{center}\\fontsize{2.8mm}{0.8em}\\selectfont \\setlength{\\arrayrulewidth}{0.7pt}\n'\
        +'\\begin{tabular}{cccc}\n'\
        +'&&&\\\\[-1.3cm]\n'\
        +'\\multicolumn{1}{c}{\\textbf{Dominio}} & \\multicolumn{1}{c}{\\Bold{Departamento}}&\\multicolumn{1}{c}{\\Bold{UPMS}}&\\multicolumn{1}{c}{\\Bold{Hogares}}\\\\[0.05cm]\\hline \n'\
        +self.tabla\
        +'&\\textbf{Total}	&\\textbf{1,037}&	\\textbf{11,540}\\\\[0.05cm]  \n'\
        +'\\hline \n'\
        +'&&&\\\\[-0.36cm]\n'\
        +'\\end{tabular}\n'\
        +'\\end{center}\n'\
        +'{\\footnotesize Fuente:  Encuesta Nacional de Condiciones de Vida (Encovi) 2014.} }' )

    def crear_presentacion_pp(self):
        print self.ruta_salida,self.titulo_documento + '-Presentacion-plantilla.pptx'
        self.prs = Presentation(os.path.join( self.ruta_salida,'Presentacion-plantilla.pptx' ))
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = 'Resultados ENCOVI 2014'
        subtitle = slide.placeholders[1]
        subtitle.text = 'Participación ciudadana y medios de comunicación \n ' + self.lugar_geografico+', noviembre 2016\n'
        title_slide_layout = self.prs.slide_layouts[3]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = 'Objetivo General de Encovi'
        subtitle = slide.placeholders[10]
        subtitle.text = 'Conocer y evaluar las condiciones de vida de la población '\
        +' y determinar los niveles de pobreza existentes en Guatemala.'
        title_slide_layout = self.prs.slide_layouts[3]
        slide = self.prs.slides.add_slide(title_slide_layout)
        shapes = slide.shapes
        title = slide.shapes.title
        title.text = 'Objetivos específicos de Encovi'
        body_shape = shapes.placeholders[10]
        tf = body_shape.text_frame
        p = tf.add_paragraph()
        p.text =  'Contar con información confiable y oportuna que permita identificar las condiciones de vida de los distintos grupos sociales del país, especialmente en la estructura de los ingresos y gastos del hogar, que faciliten la elaboración y evaluación de planes, políticas y estrategias de desarrollo.'
        p.level = 1
        p = tf.add_paragraph()
        p.text = 'Obtener estimaciones de la tasa de pobreza y pobreza extrema para cada uno de los dominios de estudio de esta encuesta. '
        p.level = 1
        p = tf.add_paragraph()
        p.text  = 'Monitorear los avances e impactos de los programas y acciones sociales. '
        p.level = 1
        for shape in slide.placeholders:
            print('%d %s' % (shape.placeholder_format.idx, shape.name))

    def terminar_presentacion_pp(self):
        self.prs.save(os.path.join( self.ruta_salida,self.titulo_documento + '-Presentacion.pptx' ))

    def crear_cajita(self, titulo,descripcion, titulo_grafica,
        des_grafica, grafica, fuente):
        cajita = '\\cajita{% \n' \
        + titulo + ' % \n' \
        + '}% \n' \
        + '{% \n' \
        + descripcion +  ' % \n' \
        + '}% \n' \
        + '{% \n'\
        + titulo_grafica + ' % \n' \
        + '}% \n' \
        + '{% \n ' \
        + des_grafica + ' % \n' \
        + '}% \n' \
        + '{% \n ' \
        + grafica +' % \n' \
        + '}% \n' \
        + '{% \n' \
        + fuente + ' % \n' \
        + '} \n'
        return cajita

    def crear_capitulo(self,nombre_cap, descripcion_cap):
        capitulo = "\\INEchaptercarta{" + nombre_cap \
        +"}{" + descripcion_cap + "} "
        return capitulo

    def crear_capitulo_presentacion(self,nombre_cap, descripcion_cap):
        capitulo = "\\INEpartecarta{" + nombre_cap \
        +"}{" + descripcion_cap + "} "
        return capitulo

    def escribir_en_doc(self, texto):
        self.documento.write( '\n \n ' + texto + '\n \n' )

    def escribir_en_presentacion(self, texto):
        self.presentacion.write( '\n \n ' + texto + '\n \n' )

    def terminar_documento(self):
        self.escribir_en_doc('\\input{metodologia.tex}\n'.encode('utf-8'))
        self.escribir_en_doc('\\end{document}'.encode('utf-8'))
        self.documento.close()


    def terminar_presentacion(self):
        self.escribir_en_presentacion('\\muchasgracias'.encode('utf-8'))
        self.escribir_en_presentacion('\\end{document}'.encode('utf-8'))
        self.presentacion.close()


    def compilar_documento(self):
        cadena_compilacion = "cd "+ self.ruta_compilacion + " && xelatex " + self.titulo_documento.strip().replace(" ", "\\ ") + ".tex"
        print cadena_compilacion
        print subprocess.Popen(cadena_compilacion, shell=True, stdout=subprocess.PIPE).stdout.read()

    def compilar_graficas(self):
        os.chdir( os.path.join(self.ruta_salida,'graficasPresentacion'))
        for file in glob.glob('*.tex'):
            cadena_compilacion = "cd "+ os.path.join(self.ruta_compilacion,'graficasPresentacion')+ " && xelatex " + file
            print cadena_compilacion
            print subprocess.Popen(cadena_compilacion, shell=True, stdout=subprocess.PIPE).stdout.read()
        os.chdir( os.path.join(self.ruta_salida,'graficasPresentacion'))
        for file in glob.glob('*.pdf'):
            cadena_compilacion = "cd "+ os.path.join(self.ruta_compilacion,'graficasPresentacion')+ " && convert -density 300 -quality 100  " + file + ' ' + os.path.splitext(file)[0] + '.png'
            print cadena_compilacion
            print subprocess.Popen(cadena_compilacion, shell=True, stdout=subprocess.PIPE).stdout.read()



    def compilar_presentacion(self):
        cadena_compilacion = "cd "+ self.ruta_compilacion + " && xelatex " + self.titulo_documento.strip().replace(" ", "\\ ") + "-Presentacion.tex"
        print cadena_compilacion
        print subprocess.Popen(cadena_compilacion, shell=True, stdout=subprocess.PIPE).stdout.read()

    def crear_cadena_descriptor(self,formato,tipo):
        archivo = open(os.path.join(self.ruta_salida, 'descripciones',formato + '.tex'), 'w')
        archivo.write('%Dummy file')
        archivo.close()
        if tipo.strip().upper() == "CUADRO":
                tex = os.path.isfile( os.path.join(self.ruta_salida, 'graficas',formato + '.tex') )
                pdf = False
        else:
            tex = os.path.isfile( os.path.join(self.ruta_salida, 'graficas',formato + '.tex') )
            pdf = os.path.isfile( os.path.join(self.ruta_salida,'graficas', formato + '.pdf') )
        retorno = ''
        if not tex or pdf:
            try:
                os.makedirs( os.path.join(self.ruta_salida,'graficas') )
            except OSError:
                print "El directorio de graficas ya existe"
            try:
                os.makedirs( os.path.join(self.ruta_salida,'cuadros') )
            except OSError:
                print "El directorio de cuadros ya existe"
            try:
                if tipo.strip().upper() == "CUADRO":
                    archivo = open(os.path.join(self.ruta_salida, 'cuadros',formato + '.tex'), 'w')
                    archivo.write('%Dummy file')
                else:
                    archivo = open(os.path.join(self.ruta_salida, 'graficas',formato + '.tex'), 'w')
                    archivo.write('%Dummy file')
            except ValueError:
                print 'Una excepcion', ValueError
                pass
        if tipo.strip().upper() == "CUADRO":
            retorno = '\\input{cuadros/' + formato + '.tex}'
        else:
            if pdf:
                retorno = '\\includepdf{' + formato + '.pdf' + '}'
            else:
                retorno = '\\begin{tikzpicture}[x=1pt,y=1pt]\\input{graficas/' + formato + '.tex}' + '\\end{tikzpicture}'
        return retorno

    def leer_csv(self, ruta):
        archivo = open(ruta, 'rb')
        lector =  csv.reader(archivo,  delimiter=';')
        salida = []
        for row in lector:
            salida.append(row)
        return salida

    def crear_carpeta_descripciones(self):
        try:
            os.makedirs( os.path.join(self.ruta_salida, 'descripciones') )
        except OSError:
            print "El directorio ya existe"

    def des_101(self):
        des = 'El departamento de ' + self.lugar_geografico + ' está ubicado '\
        + ' en la Región ' + self.datos_depto[2] + ' del país; tiene '\
        + self.datos_depto[3] + ' municipios de los cuales ' +  self.datos_depto[4] \
        + ' es la cabecera departamental. \n\n'\
        +'La superficie, en kilómetros cuadrados, del departamento es de ' + self.datos_depto[5]
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_01.tex'), 'w')
        archivo.write(des)





    def des_102(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_02.csv'))
        des = 'Saber cuál es la población de un departamento' \
        +' es fundamental para el diseño de políticas públicas, ya que este dato representa el ' \
        +'universo de los potenciales usuarios de los distintos programas y' \
        +' proyectos implementados por el Gobierno a nivel departamental. \n\n ' \
        +' Se estima que en 2014 el departamento de ' + self.lugar_geografico + ' tenía ' + self.formato_bonito(datos[3][1]) \
        +' habitantes; este dato representa un ' + self.cambio(datos[3][1], datos[2][1]) \
        +' del ' + self.porcentaje(datos[3][1], datos[2][1]) + '\\% respecto de la población estimada para 2011.' \
        + ' Para el 2006, la Encovi estimó que la población de este departamento ascendía a '+  self.formato_bonito(datos[1][1])
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_02.tex'), 'w')
        archivo.write(des)

    def des_103(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_03.csv'))
        des = 'Al comparar el número de habitantes del departamento con la población total, resulta que para 2014 el departamento de '\
        + self.lugar_geografico +  ' representaba el ' +self.formato_bonito(datos[3][1])\
        + ' de la población de la República. \n\n'\
        + ' En 2006, el departamento representaba el ' + self.formato_bonito(datos[1][1])\
        + ' de la población total.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_03.tex'), 'w')
        archivo.write(des)

    def des_104(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_04.csv'))
        des = 'La densidad poblacional relaciona el número de habitantes con el tamaño del territorio donde esta población habita. Normalmente este indicador se expresa en personas por kilómetro cuadrado. \n \n'\
        +' En 2014 la densidad poblacional de ' + self.lugar_geografico +  ' era de ' \
        +self.formato_bonito(datos[3][1]) + ' habitantes por kilómetro cuadrado; este dato es ' \
        +self.mayor_menor(datos[3][1],datos[1][1]) + ' que el estimado para 2006, el cual ascendía a  ' \
        + self.formato_bonito(datos[1][1]) + ' habitantes por kilómetro cuadrado.'

        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_04.tex'), 'w')
        archivo.write(des)

    def des_105(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_05.csv'))
        des = ' El departamento de ' + self.lugar_geografico + ' está relativamente '\
        + self.mas_menos(datos[2][1],datos[1][1]) + ' poblado que el territorio nacional\n\n'\
        + 'Efectivamente, para 2014 la densidad poblacional del departamento era de '\
        + self.formato_bonito(datos[2][1]) + ' dato ' + self.mayor_menor(datos[2][1],datos[1][1])\
        + ' que la densidad a nivel nacional, la que se ubicó en ' + self.formato_bonito(datos[1][1])

        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_05.tex'), 'w')
        archivo.write(des)

    def des_106(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_06.csv'))
        des = 'Para el 2014, la Encovi estima que el ' + self.formato_bonito(datos[2][2]) \
        +'\\% de la población del departamento de ' + self.lugar_geografico \
        + ' era mujer y el restante ' + self.formato_bonito(datos[2][1]) + '\\% hombre. \n \n' \
        + 'El departamento muestra una proporción un poco '+  self.mayor_menor(datos[2][2],datos[1][2])\
        +'de población femenina que el indicador nacional.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_06.tex'), 'w')
        archivo.write(des)

    def des_107(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_07.csv'))
        des = 'Debido a que las personas de distintas edades tienen diferentes necesidades, es importante conocer la estructura de la población por grupos de edad. Por ejemplo, es de mucho interés saber cuántos niños hay en un departamento porque ellos necesitan educarse para su adecuado desarrollo; el número de adultos mayores también es relevante, porque este grupo necesita de mayores cuidados médicos. \n\n'\
        +'La Encovi estima que en 2014 el ' + self.formato_bonito(datos[1][1]) +  '\\% de la población del departamento de '\
        + self.lugar_geografico + ' es menor de quince años, mientras que el '\
        + self.formato_bonito(datos[4][1]) + '\\% tiene 65 años o más.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_07.tex'), 'w')
        archivo.write(des)

    def des_108(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_08.csv'))
        des = 'Debido a que la población menor de edad requiere de una atención especial por parte del Gobierno, es necesario conocer su magnitud para cada departamento.\n\n'\
        +' Para 2014 en el departamento de ' + self.lugar_geografico + ' el '\
        +self.formato_bonito(datos[2][1]) + ' \\% de su población era '\
        +' que el dato nacional de ' + self.formato_bonito(datos[1][1]) + '\\%'

        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_08.tex'), 'w')
        archivo.write(des)

    def des_109(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_09.csv'))
        des = 'Las personas mayores de 65 años también son más vulnerables debido a que en mucho casos ya no generan ingresos y porque están en mayor riesgo de contraer enfermedades.\n\n'\
        +'Según los resultados de la Encovi 2014, la proporción de población mayor de 65 años en el departamento de '\
        +self.lugar_geografico + ' era de ' + self.formato_bonito(datos[2][1])\
        +'\\%, porcentaje ' + self.mayor_menor(datos[2][1],datos[1][1])\
        + ' que el dato nacional de ' + self.formato_bonito(datos[1][1]) + '\\%'

        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_09.tex'), 'w')
        archivo.write(des)

    def des_110(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_10.csv'))
        des = 'Una de las mayores riquezas de Guatemala es su diversidad cultural, la cual puede analizarse más a detalle en los departamentos. En el 2006, el '\
        +self.formato_bonito(datos[1][1]) + '\\% de la población del departamento'\
        +' se autoidentificaba como indígena (Maya, Xinca o Garífuna) \n\n'\
        +' Para 2014 este indicador se ubicó en ' + self.formato_bonito(datos[3][1]) + '\\%'


        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_10.tex'), 'w')
        archivo.write(des)

    def des_111(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_11.csv'))
        des = 'Guatemala es un país pluricultural, en este sentido la Encovi 2014 muestra que el 38.8\\% de la población a nivel nacional se autoindentifica como perteneciente al pueblo Maya, Xinca o Garífuna. \n\n'\
        +'Este indicador para el departamento de ' + self.lugar_geografico\
        +' es ' + self.mayor_menor(datos[2][1],datos[1][1]) + ' que el dato nacional'\
        +', al ascender a ' + self.formato_bonito(datos[2][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_11.tex'), 'w')
        archivo.write(des)

    def des_112(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_12.csv'))
        des = 'Guatemala es un país multilingüe. A nivel nacional un 30\\% de la población aprendió a hablar en un idioma distinto al español, según los datos de la Encovi 2014.\n\n'\
        +' En el departamento de ' + self.lugar_geografico + ', el '\
        +self.formato_bonito(datos[1][1]) + '\\% de la población aprendió a '\
        + 'hablar en un idioma maya y el ' + self.formato_bonito(datos[2][1])\
        + ' restante su idioma materno es el español.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_12.tex'), 'w')
        archivo.write(des)


    def des_113(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_13.csv'))
        des = 'Debido a que el País es multilingüe, la población que habla un idioma diferente al materno es importante, por la capacidad de comunicación que genera entre los diferentes grupos lingüísticos.\n\n'\
        +'En el 2014, en el departamento de ' + self.lugar_geografico + ' el '\
        +self.formato_bonito(datos[3][1]) + '\\%' + ' de la población hablaba '\
        +' más de un idioma. En el 2006 este indicador ascendía a ' +self.formato_bonito(datos[1][1])
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_13.tex'), 'w')
        archivo.write(des)

    def des_114(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_14.csv'))
        des = 'La Encovi 2014 muestra que en el departamento de ' + self.lugar_geografico\
        +', el porcentaje de personas indígenas que hablan más de un idioma es de '\
        +self.formato_bonito(datos[1][1]) + '\\%, mientas que esta proporción para '\
        +' las personas no indígenas es de ' + self.formato_bonito(datos[2][1])\
        +'\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_14.tex'), 'w')
        archivo.write(des)

    def des_115(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_15.csv'))
        des = 'De la población que habla más de un idioma en el departamento de '\
        + self.lugar_geografico+ ', el '+  self.formato_bonito(datos[1][1])\
        +'\\% '+'habla un idioma indígena como segundo idioma, el  '+ self.formato_bonito(datos[2][1])\
        +'\\% ' + ' el español es su segunda lengua y el ' + self.formato_bonito(datos[3][1])\
        +'\\% habla otro idioma.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_15.tex'), 'w')
        archivo.write(des)

    def des_116(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_16.csv'))
        des = 'En el departamento de ' + self.lugar_geografico + ' la proporción'\
        +' de la población que solo habla el idioma predominante fue de '\
        +self.formato_bonito(datos[3][1]) + '\\% en el 2014. \n\n'\
        +' Este indicador se ubicó en ' + self.formato_bonito(datos[1][1]) + '\\%'\
        +' en 2006.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_16.tex'), 'w')
        archivo.write(des)

    def des_117(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_17.csv'))
        des = 'En el departamento de ' + self.lugar_geografico + ' la proporción '\
        +' de la población que habla el idioma predominante y el español fue de '\
        +self.formato_bonito(datos[3][1]) + '\\% en el 2014.\n\n'\
        +' En indicador se ubicó en ' + self.formato_bonito(datos[1][1]) + '\\%'\
        +' en 2006.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_17.tex'), 'w')
        archivo.write(des)

    def des_118(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_18.csv'))
        des = 'La Encovi 2014 establece que en 2014 el 49.5\\% de la población'\
        +' a nivel nacional habita en áreas urbanas.\n\n'\
        +' En el departamento de ' + self.lugar_geografico + ' este porcentaje es '\
        +self.mayor_menor(datos[2][1],datos[1][1]) + ' que el dato nacional al '\
        +' ubicarse en ' + self.formato_bonito(datos[2][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_18.tex'), 'w')
        archivo.write(des)

    def des_119(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_19.csv'))
        des = 'Es interesante conocer la proporción que habita en áreas urbanas atendiendo al grupo étnico.\n\n'\
        +' Según la Encovi 2014, el ' + self.formato_bonito(datos[1][1]) + '\\% '\
        +' de la población indígena del departamento de ' + self.lugar_geografico\
        +' habita en áreas urbanas; para la población no indígena la proporción '\
        +' de habitantes en áreas urbanas asciende a ' + self.formato_bonito(datos[2][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_19.tex'), 'w')
        archivo.write(des)

    def des_120(self):
        datos1 = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_18.csv'))
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_20.csv'))
        des = 'Como se mencionó anteriormente, en el departamento de '\
        + self.lugar_geografico + ' el ' + self.formato_bonito(datos1[2][1])\
        +'\\% de la población habita en áreas urbanas. \n\n'\
        +'Este dato es ' + self.mayor_menor(datos[3][1],datos[1][1]) + ' que el '\
        +'observado en la Encovi de 2006, en la que el indicador se ubicó en '\
        +self.formato_bonito(datos[1][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_20.tex'), 'w')
        archivo.write(des)


    def des_121(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_21.csv'))
        des = 'El Documento Personal de Identificación (DPI) es necesario para llevar a cabo una serie de gestiones en el ámbito personal, comercial e institucional. A nivel nacional el 96.4\\% de la población mayor de edad tenía DPI en el 2014.\n\n'\
        +'Para el departamento de ' + self.lugar_geografico + ' este porcentaje se ubicó en '\
        +self.formato_bonito(datos[2][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_21.tex'), 'w')
        archivo.write(des)

    def des_122(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_22.csv'))
        des = 'Por sexo, el ' + self.formato_bonito(datos[1][1]) + '\\% de los hombres '\
        +' del departamento de ' + self.lugar_geografico + ' tenían DPI en el 2014, '\
        +'dato ligeramente ' + self.mayor_menor(datos[1][1],datos[2][1]) + ' que el'\
        +' de las mujeres que se ubicó en ' +  self.formato_bonito(datos[2][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_22.tex'), 'w')
        archivo.write(des)


    def des_123(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','1_23.csv'))
        des = 'En cuanto al área de residencia, el ' + self.formato_bonito(datos[1][1])\
        +'\\% de los habitantes mayores de edad de las áreas urbanas del departamento '\
        +' de ' + self.lugar_geografico + ' tenían DPI en el 2014, dato '\
        +self.mayor_menor(datos[1][1],datos[2][1]) + ' al de las áreas rurales '\
        +' que se ubicó en ' + self.formato_bonito(datos[2][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','1_23.tex'), 'w')
        archivo.write(des)

    def des_201(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_01.csv'))
        des = 'El tipo de material del piso de una vivienda incide en la calidad de vida de sus habitantes. Por ejemplo, el piso de tierra es de difícil limpieza, por lo que las personas que habitan viviendas con este tipo de material tienen mayor riesgo de contraer enfermedades gastrointestinales. \n\n' \
        +'En este sentido, la Encovi 2014 muestra que el ' + self.formato_bonito(datos[7][1]) \
        +'\\% de hogares en el departamento de ' + self.lugar_geografico \
        +' habitan viviendas con piso de tierra, mientras que el ' \
        + self.formato_bonito(datos[4][1]) +'\\% habita en viviendas con piso de torta de cemento.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_01.tex'), 'w')
        archivo.write(des)


    def des_202(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_02.csv'))
        des = 'La Encovi 2014 establece que, a nivel nacional, el 27.5\\% de hogares habita en una vivienda con piso de tierra.\n\n'\
        +'Este porcentaje para el departamento de ' + self.lugar_geografico\
        +' es ' + self.mayor_menor(datos[2][1],datos[1][1]) + ' que el dato '\
        +'a nivel nacional y se ubica en ' + self.formato_bonito(datos[2][1])
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_02.tex'), 'w')
        archivo.write(des)



    def des_203(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_03.csv'))
        des = 'El material del techo también es un aspecto relevante para determinar las condiciones de vida de un hogar. Un techo de mala calidad no brinda el suficiente resguardo para las inclemencias del tiempo, lo cual puede incidir en una mayor prevalencia de enfermedades respiratorias. \n\n' \
        +'En el departamento de ' + self.lugar_geografico \
        +' el ' + self.formato_bonito(datos[2][1]) + '\\% de hogares habitan viviendas con techo de lámina, ' + self.especial(datos[5][1])
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_03.tex'), 'w')
        archivo.write(des)


    def des_204(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_04.csv'))
        des = 'A nivel nacional, el 72.3\\% de hogares habita en una vivienda con techo de lámina,'\
        +' según la Encovi 2014.\n\n'\
        +'Para el caso del departamento de ' + self.lugar_geografico + ', el porcentaje '\
        +' de hogares con este tipo de material en el techo es ' + self.mayor_menor(datos[2][1],datos[1][1])\
        +' que el valor del dato nacional, al ubicarse en ' + self.formato_bonito(datos[2][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_04.tex'), 'w')
        archivo.write(des)



    def des_205(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_05.csv'))
        des = 'Se sabe que hay algunas enfermedades cuyos vectores se reproducen más fácilmente en cierto tipo de paredes, como el mal de Chagas; por ello es relevante investigar acerca del tipo de material de las paredes en el que habitan los hogares. \n\n'\
        +'En el 2014, el ' + self.formato_bonito(datos[2][1])+ '\\% de hogares del departamento de '\
        +self.lugar_geografico + ' habitaban una casa con paredes de block, '\
        +self.especial1(datos[4][1])
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_05.tex'), 'w')
        archivo.write(des)


    def des_206(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_06.csv'))
        des = 'Los hogares que habitan en viviendas con paredes de block representaban el 58.2\\% del total de hogares en 2014.\n\n'\
        +'En el departamento de ' + self.lugar_geografico + ', el porcentaje de hogares que '\
        +' habitan viviendas con paredes de block ascendía a ' + self.formato_bonito(datos[2][1])\
        +'\\% para ese mismo año.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_06.tex'), 'w')
        archivo.write(des)


    def des_207(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_07.csv'))
        des = 'Tener acceso a una de red de agua es muy importante para la salud de los miembros de un hogar.\n\n'\
        +' En el 2014, el ' + self.formato_bonito(datos[3][1]) +'\\% '\
        +' de hogares del departamento de ' + self.lugar_geografico + ' estaban '\
        +' conectados a una red de agua; en 2006 este porcentaje era de '\
        +self.formato_bonito(datos[1][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_07.tex'), 'w')
        archivo.write(des)

    def des_208(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_08.csv'))
        des = 'Al igual que la red de agua, el acceso a una red de drenajes también es esencial para la higiene y saneamiento de un hogar.\n\n'\
        +' Según la Encovi 2014, el ' + self.formato_bonito(datos[3][1])\
        +'\\% de hogares del departamento de ' + self.lugar_geografico\
        +' estaban conectados a una red de drenajes; en 2006 este porcentaje '\
        +'ascendía a ' + self.formato_bonito(datos[1][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_08.tex'), 'w')
        archivo.write(des)

    def des_209(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_09.csv'))
        des = 'Para el caso del acceso a una red de energía eléctrica, la Encovi 2014'\
        +' muestra que el ' + self.formato_bonito(datos[3][1]) + '\\% de los hogares '\
        +' del departamento de ' + self.lugar_geografico + ' estaban conectados a '\
        +' este tipo de red. \n\n ' + ' El porcentaje para el año 2006 se ubicaba en '\
        + self.formato_bonito(datos[1][1]) + '\\%, según la misma fuente de información.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_09.tex'), 'w')
        archivo.write(des)

    def des_210(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_10.csv'))
        des = 'En el departamento de ' + self.lugar_geografico + ' el porcentaje '\
        +' de hogares que tenían automóvil ascendía a ' + self.formato_bonito(datos[2][1])\
        +'\\%, porcentaje ' + self.mayor_menor(datos[2][1], datos[1][1]) + ' al '\
        +'reportado a nivel nacional. '
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_10.tex'), 'w')
        archivo.write(des)

    def des_211(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_11.csv'))
        des = 'En cuanto a los hogares con televisión, en  el departamento de ' + self.lugar_geografico + ' el porcentaje '\
        +' el porcentaje de hogares que contaban con este artículo era de  ' + self.formato_bonito(datos[2][1])\
        +'\\%, indicador ' + self.mayor_menor(datos[2][1], datos[1][1]) + ' al '\
        +'reportado a nivel nacional. '
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_11.tex'), 'w')
        archivo.write(des)

    def des_212(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_12.csv'))
        des = 'En el departamento de ' + self.lugar_geografico + ' el porcentaje '\
        +' de hogares que contaban con refrigeradora era de  ' + self.formato_bonito(datos[2][1])\
        +'\\%, porcentaje ' + self.mayor_menor(datos[2][1], datos[1][1]) + ' al '\
        +'reportado a nivel nacional. '
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_12.tex'), 'w')
        archivo.write(des)

    def des_213(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_13.csv'))
        des = 'En el año 2006, el  ' + self.formato_bonito(datos[3][1]) + ' \\% '\
        +' de hogares en el departamento  de ' + self.lugar_geografico + ' eran '\
        +' nucleares, es decir, estaban compuestos por padres e hijos únicamente; '\
        ' en el otro extremo el ' + self.formato_bonito(datos[1][1]) + '\\% de '\
        +' hogares tenían un único miembro (unipersonales).\n\n'\
        +' Para el caso de los hogares conformados solamente por la pareja (sin hijos)'\
        + ', estos representaban el ' + self.formato_bonito(datos[2][1]) + '\\%.'

        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_13.tex'), 'w')
        archivo.write(des)

    def des_214(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_14.csv'))
        des = 'La Encovi 2014 establece que el '+ self.formato_bonito(datos[3][1])\
        +'\\% de hogares en el departamento de ' + self.lugar_geografico + ' eran nucleares, el '\
        + self.formato_bonito(datos[1][1]) + '\\% tenían un único miembro y el '\
        +self.formato_bonito(datos[2][1])+ '\\% estaban conformados por una pareja.\n\n'\
        +'Para este año, los hogares biparentales extensos (que además de los hijos tenían otros parientes viviendo en el hogar) '\
        +' representaban el ' + self.formato_bonito(datos[4][1]) + '\\%.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_14.tex'), 'w')
        archivo.write(des)


    def des_215(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_15.csv'))
        des = 'En el departamento de ' + self.lugar_geografico + ' el porcentaje '\
        +' de hogares con jefatura femenia e hijos era de ' + self.formato_bonito(datos[2][1])\
        + '\\%. \n\n' + ' Este indicador en 2006 ascendía a ' + self.formato_bonito(datos[1][1])\
        + '\\% según la Encovi de ese año.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_15.tex'), 'w')
        archivo.write(des)

    def des_216(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_16.csv'))
        des = 'En el 2014 en el departamento de '+ self.lugar_geografico + ' el '\
        +' porcentaje de hogares con jefatura femenina y en los que también habitaban '\
        +' parientes era de ' + self.formato_bonito(datos[2][1]) + '\\%.\n\n'\
        +'Este indicador en 2006 ascendía a ' + self.formato_bonito(datos[1][1])\
        +'\\%, según la Encovi de ese año.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_16.tex'), 'w')
        archivo.write(des)


    def des_217(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_17.csv'))
        des = 'El tamaño de un hogar es una variable demográfica relevante. Por ejemplo, este indicador es útil para el diseño de soluciones habitacionales que se ajusten a las necesidades de las familias.\n\n'\
        +'Según los datos de la Encovi, en 2014 los hogares del departamento  de '\
        + self.lugar_geografico + ' tenían en promedio ' + self.formato_bonito(datos[3][1])\
        + ' miembros, dato ' + self.mayor_menor(datos[3][1], datos[1][1])\
        +' a los ' + self.formato_bonito(datos[1][1]) + ' habitantes promedio por hogar del año 2006.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_17.tex'), 'w')
        archivo.write(des)

    def des_218(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_18.csv'))
        des = 'Según área de residencia, la Encovi 2014  muestra que los hogares del área rural del departamento de '\
        +self.lugar_geografico + ' son un poco ' + self.plural_mayor_menor(datos[1][1],datos[2][1]) \
        +' que los del área urbana. \n\n '\
        +'Concretamente, esta encuesta estima que en este departamento el hogar promedio rural tiene '\
        +self.formato_bonito(datos[2][1]) + ' miembros y el urbano ' + self.formato_bonito(datos[1][1])
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_18.tex'), 'w')
        archivo.write(des)

    def des_219(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','2_19.csv'))
        des = 'La Encovi revela que a nivel nacional los hogares pobres tienen, en promedio, más habitantes que los no pobres. \n \n'\
        +'Esta tendencia también se encuentra en el departamento de ' + self.lugar_geografico \
        + ', donde en el 2014 los hogares pobres extremos tenían en promedio '\
        +self.formato_bonito(datos[1][1]) + ' habitantes, dato mayor a los '\
        +self.formato_bonito(datos[3][1]) + ' miembros de los hogares no pobres. '
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','2_19.tex'), 'w')
        archivo.write(des)

    def des_601(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_01.csv'))
        des = 'Para el año 2006, el valor de la línea de pobreza total era de Q 6,574. Es importante recordar que la línea de pobreza total incluye, además del costo de una canasta básica de alimentos, un monto adicional que corresponde al porcentaje de consumo no alimenticio de las personas, cuyo consumo de alimentos se encuentra alrededor de la línea de pobreza extrema. \n \n'\
        +'Se puede observar que para 2014, el valor de la línea de pobreza total aumentó a Q 10,218, que equivale a un incremento del 137\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_01.tex'), 'w')
        archivo.write(des)

    def des_602(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_02.csv'))
        des ='Al comparar el consumo de las familias con la línea de pobreza total, resulta que en 2014 el '\
        +self.formato_bonito(datos[2][1]) + ' \\% de personas en el departamento de '\
        +self.lugar_geografico + ' se encontraba en condición de pobreza. \n \n'\
        +' Este porcentaje es más ' + self.alto_bajo(datos[2][1], datos[1][1]) + '  que el observado en 2006, el cual ascendía a '\
        +self.formato_bonito(datos[1][1])  +'\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_02.tex'), 'w')
        archivo.write(des)

    def des_603(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_03.csv'))
        des = 'La Encovi 2014 revela que el 59.3\\% de la población guatemalteca total se encontraba en condición de pobreza. Sin embargo, es importante analizar este dato territorialmente para determinar diferencias a lo interno del país. \n\n'\
        +'Al comparar la incidencia de pobreza del departamento de ' + self.lugar_geografico\
        +' con el dato nacional, se observa que éste tiene un porcentaje de pobreza total '\
        +self.mayor_menor(datos[2][1],datos[1][1]) + ' que el promedio de todo el país. '
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_03.tex'), 'w')
        archivo.write(des)

    def des_604(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_04.csv'))
        des = ' Por área de residencia, la Encovi 2014  muestra que los hogares del área rural del departamento de '\
        +self.lugar_geografico + ' tienen una  ' + self.mayor_menor(datos[2][1], datos[1][1])\
        +' incidencia de pobreza que los del área urbana. \n\n '\
        + ' Concretamente, esta encuesta estima que en este departamento la incidencia de la pobreza total en el área rural es del  '\
        +self.formato_bonito(datos[2][1]) + '\\%,  mientras que el área urbana es de  ' + self.formato_bonito(datos[1][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_04.tex'), 'w')
        archivo.write(des)

    def des_605(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_05.csv'))
        des = ' En el departamento de ' + self.lugar_geografico + ' el '+ self.formato_bonito(datos[1][1]) +'\\% '\
        +' de la población que se autoidentificaba como indígena tenía un consumo '\
        +'por debajo de la línea de pobreza total. \n\n '\
        +'Para el caso de la población no indígena, el porecentaje de pobreza '\
        +' total fue de '+ self.formato_bonito(datos[2][1]) + '\\% en 2014.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_05.tex'), 'w')
        archivo.write(des)

    def des_606(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_06.csv'))
        des = ' A nivel nacional, el 68.2\\% de la población menor de edad se encontraba en nivel de pobreza en el año 2014, según los dato de la Encovi.\n\n'\
        +' Para el departamento de ' + self.lugar_geografico + ', este indicador está por '\
        + self.encima_debajo(datos[2][1],datos[1][1]) + ' del dato nacional ya que '\
        +' asciende a ' + self.formato_bonito(datos[2][1]) + '\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_06.tex'), 'w')
        archivo.write(des)


    def des_607(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_07.csv'))
        des = ' Para el año 2006, el valor de la línea de pobreza extrema era de Q 3,206. Es importante recordar que la línea de pobreza extrema representa el costo de adquirir la cantidad de calorías mínimas recomendadas para un humano. \n\n'\
        +'Se puede observar que para 2014, el valor de la línea de pobreza extrema aumentó a Q 5,750 que equivale a un incremento del 79.4% respecto al valor de 2006.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_07.tex'), 'w')
        archivo.write(des)


    def des_608(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_08.csv'))
        des = 'Al comparar el consumo de las familias con la línea de pobreza extrema, resulta que en 2014 el '\
        +self.formato_bonito(datos[2][1]) + '\\% de personas en el departamento de '\
        +self.lugar_geografico + ' se encontraba en condición de pobreza extrema.\n\n'\
        +'Este porcentaje es más ' + self.alto_bajo(datos[2][1],datos[1][1]) + ' que el observado en 2006, el cual ascendía a '\
        +self.formato_bonito(datos[1][1]) +'\\%'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_08.tex'), 'w')
        archivo.write(des)

    def des_609(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_09.csv'))
        des = 'En 2014, el 23.4\\% de la población guatemalteca se encontraba en condición de pobreza extrema. Gracias al diseño muestral de la Encovi, es posible desagregar este dato por departamento. \n \n'\
        +'La incidencia de pobreza extrema en el departamento de ' + self.lugar_geografico \
        +' en el 2014 fue de ' + self.formato_bonito(datos[2][1]) + ', dato '\
        +self.mayor_menor(datos[2][1],datos[1][1]) + ' que el porcentaje nacional.'

        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_09.tex'), 'w')
        archivo.write(des)

    def des_610(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_10.csv'))
        des = 'En 2014, el 11.2\\% de la población urbana estaba en pobreza extrema; para el caso de la población rural, este indicador se ubicó en 35.3\\%.\n\n'\
        +'Para el departamento de ' + self.lugar_geografico \
        +', la pobreza extrema en el área rural se ubicó en ' + self.formato_bonito(datos[2][1]) \
        +'\\% y en el área urbana en ' + self.formato_bonito(datos[1][1]) + '\\%, según la información de la Encovi 2014.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_10.tex'), 'w')
        archivo.write(des)

    def des_611(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_11.csv'))
        des = 'El ' + self.formato_bonito(datos[1][1]) + '\\% de la población '\
        +' indígena del departamento de '+ self.lugar_geografico + ' estaba con '\
        +' condición de pobreza extrema en el 2014.\n\n'\
        +'Para el caso del la población no indígena, la indicencia de pobreza extrema '\
        +' fue de ' + self.formato_bonito(datos[2][1])+ '\\% para ese mismo año.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_11.tex'), 'w')
        archivo.write(des)

    def des_612(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_12.csv'))
        des = 'A nivel nacional, el 29.0\\% de la población menor de edad se encontraba en nivel de pobreza extrema en el año 2014, según los dato de la Encovi.\n\n'\
        +'Para el departamento de ' + self.lugar_geografico +', este indicador está por '\
        +self.encima_debajo(datos[2][1], datos[1][1]) + ' del dato nacional, al ubicarse en '\
        +self.formato_bonito(datos[2][1]) + '\\%.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_12.tex'), 'w')
        archivo.write(des)

    def des_613(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','6_13.csv'))
        des = 'A nivel nacional, el índice de Gini para medir la desigualdad se ubicó en 0.53 en el 2014 según los datos de la Encovi.\n\n'\
        +' Para el caso del departamento de ' + self.lugar_geografico + ', este indicador '\
        +' estuvo por ' + self.encima_debajo(datos[2][1],datos[1][1])\
        +' del dato nacional al ubicarse en ' + self.formato_bonito(datos[2][1])
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','6_13.tex'), 'w')
        archivo.write(des)


    def des_301(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','3_01.csv'))
        des = 'Contar con un seguro de salud garantiza, en buena medida, el acceso a servicios médicos para el tratamiento de enfermedades y lesiones ocasionadas por accidentes.\n\n'\
        +'Para el caso del departamento de ' + self.lugar_geografico \
        +', la Encovi 2014 revela que el ' + self.formato_bonito(datos[3][1])\
        +'\\% de la población no estaba cubierta por ningún tipo de seguro médico, ni privado ni público.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','3_01.tex'), 'w')
        archivo.write(des)

    def des_302(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','3_02.csv'))
        des = 'El seguro social fue una de las principales conquistas de la Revolución de octubre de 1944, el cual fue creado con el objetivo de ser una garantía de salud para todos los guatemaltecos. \n\n'\
        +'La Encovi 2014 muestra que el ' + self.formato_bonito(datos[3][1]) \
        +'\\% de la población del departamento de ' + self.lugar_geografico\
        +' tuvo acceso a esta protección social. '
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','3_02.tex'), 'w')
        archivo.write(des)

    def des_303(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','3_03.csv'))
        des = 'El papanicolau es un examen médico que tiene como objetivo el diagnóstico preventivo del cáncer del cuello uterino. \n\n'\
        +'Según la información de la Encovi, en 2014 el '+ self.formato_bonito(datos[3][1])\
        +'\\% de las mujeres en edad fértil del departamento de ' + self.lugar_geografico\
        +' se realizaron alguna vez el examen de papanicolau.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','3_03.tex'), 'w')
        archivo.write(des)

    def des_304(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','3_04.csv'))
        des = 'Las mamografías son exámenes que buscan detectar en forma temprana el cáncer de seno. A nivel nacional, en 2014 el 4.2\\% de mujeres entre 15 a 49 años se habían realizado este tipo de examen en los doce meses anteriores a la encuesta.\n\n'\
        +'Para las mujeres en edad fértil del departamento de ' + self.lugar_geografico \
        +', este porcentaje se ubicó en ' + self.formato_bonito(datos[3][1])+ '\\%, según la Encovi 2014.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','3_04.tex'), 'w')
        archivo.write(des)

    def des_305(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','3_05.csv'))
        des = 'El promedio de embarazos de las mujeres en edad fértil es una variable importante tanto para la salud materno infantil, como para el estudio de las tendencias demográficas de un país.\n\n'\
        +'La Encovi 2014 muestra que el departamento de ' + self.lugar_geografico \
        +' las mujeres en edad fértil han tenido, en promedio, ' + self.formato_bonito(datos[3][1]) + ' embarazos.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','3_05.tex'), 'w')
        archivo.write(des)

    def des_401(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','4_01.csv'))
        des = 'La alfabetización universal es uno de los objetivos de desarrollo más importantes del país. A nivel nacional, el 79.1\\% de la población de 15 años o más sabía leer y escribir en el 2014, según los datos de la Encovi.\n\n'\
        +'En el departamento de ' + self.lugar_geografico + ' , el '\
        +self.formato_bonito(datos[3][1]) +  '\\% de los mayores de 14 años sabían leer y escribir en 2014 según revela la información de la Encovi de ese año. Se observa en la gráfica que este indicador ha tenido una tendencia creciente.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','4_01.tex'), 'w')
        archivo.write(des)

    def des_402(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','4_02.csv'))
        des = 'Por pobreza, la Encovi 2014 muestra que la tasa de alfabetismo en el departamento de '\
        +self.lugar_geografico + ' para los pobres extremos era de '\
        +self.formato_bonito(datos[1][1]) + '\\%, para los pobres no extremos de ' + self.formato_bonito(datos[2][1])\
        +'\\%. En general puede observarse que a mayor pobreza, menor es la tasa de alfabetismo.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','4_02.tex'), 'w')
        archivo.write(des)

    def des_403(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','4_03.csv'))
        des = 'Los años de escolaridad promedio miden que tanto ha avanzado una población en los distintos niveles y grados educativos. \n\n'\
        +'En el departamento de ' + self.lugar_geografico + ' , la Encovi de 2014 señala que, en promedio, la población alcanzó '\
        +self.formato_bonito(datos[3][1]) +' años de escolaridad.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','4_03.tex'), 'w')
        archivo.write(des)

    def des_404(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','4_04.csv'))
        des = 'Que los niños no asistan a la escuela es una condición no deseada, debido a la importancia que la educación tiene en los niños para su correcto desarrollo.\n\n'\
        +'La Encovi 2014 muestra que en el departamento de ' + self.lugar_geografico \
        +', el ' + self.formato_bonito(datos[3][1]) + '\\% de niños entre 7 y 12 años no estaban inscritos en un centro educativo del nivel primario.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','4_04.tex'), 'w')
        archivo.write(des)

    def des_405(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','4_05.csv'))
        des = 'En general, los países que logran que buena parte de población joven ingrese a la educación media, tienen mejores niveles de desarrollo.\n\n'\
        +'Para el caso del departamento de ' + self.lugar_geografico + ', el '\
        +self.formato_bonito(datos[3][1])+ '\\% no logró asistir a un plantel educativo del ciclo básico en 2014, según la información de la Encovi.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','4_05.tex'), 'w')
        archivo.write(des)

    def des_405(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','4_05.csv'))
        des = 'En general, los países que logran que buena parte de población joven ingrese a la educación media, tienen mejores niveles de desarrollo.\n\n'\
        +'Para el caso del departamento de ' + self.lugar_geografico + ', el '\
        +self.formato_bonito(datos[3][1])+ '\\% no logró asistir a un plantel educativo del ciclo básico en 2014, según la información de la Encovi.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','4_05.tex'), 'w')
        archivo.write(des)

    def des_501(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','5_01.csv'))
        des = 'La tasa de participación –la proporción de la población económicamente activa respecto de la población en edad de trabajar- en el departamento de '\
        + self.lugar_geografico + ' fue de ' + self.formato_bonito(datos[3][1])\
        +'\\% en el año 2014. \n\n'\
        +'Este indicador en el año 2006 se situó en ' + self.formato_bonito(datos[1][1])\
        +'\\%, según datos de la Encovi de ese año.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','5_01.tex'), 'w')
        archivo.write(des)

    def des_502(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','5_02.csv'))
        des = 'Si se desagrega por sexo, la tasa de participación de los hombres es '\
        + self.mayor_menor(datos[1][1],datos[2][1]) + ' que la de las mujeres. \n\n'\
        +'Efectivamente, mientras que para los hombres la tasa de participación se situó '\
        +'en ' + self.formato_bonito(datos[1][1]) + '\\%, para las mujeres el indicador '\
        +' se ubicó en ' + self.formato_bonito(datos[2][1]) + '\\%.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','5_02.tex'), 'w')
        archivo.write(des)

    def des_503(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','5_03.csv'))
        des = 'Por área de residencia la Encovi 2014 muestra que la tasa de participación es '\
        + self.mayor_menor(datos[2][1],datos[1][1]) + ' en las áreas rurales que en las urbanas. \n\n'\
        +'En efecto, este indicador es de ' + self.formato_bonito(datos[2][1])\
        +'\\% en el área rural y de '+ self.formato_bonito(datos[1][1])\
        +'\\% en la urbana.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','5_03.tex'), 'w')
        archivo.write(des)


    def des_504(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','5_04.csv'))
        des = 'Según la legislación nacional y los convenios internacionales en materia de trabajo, no es permitido que los menores de quince años trabajen. \n\n'\
        +'La Encovi 2014 muestra que en el departamento de ' + self.lugar_geografico\
        +' el ' + self.formato_bonito(datos[3][1]) + '\\% de niños entre 7 y 14 años llevaban '\
        +' a cabo una actividad económica; este dato es ' + self.mayor_menor(datos[3][1],datos[1][1])\
        +' al observado en la encuesta de 2006.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','5_04.tex'), 'w')
        archivo.write(des)

    def des_505(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','5_05.csv'))
        des = 'En el departamento de ' + self.lugar_geografico + '  el '\
        +self.formato_bonito(datos[1][1]) + '\\% de los ocupados labora en el sector primario de la economía (agricultura, silvicultura, pesca, etc.), el '\
        +self.formato_bonito(datos[2][1]) + '\\% en la industria y el '\
        +self.formato_bonito(datos[3][1]) + '\\% en los servicios.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','5_05.tex'), 'w')
        archivo.write(des)

    def des_506(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','5_07.csv'))
        des = 'A nivel nacional el 29.1 de la población ocupada  no recibe ninguna remuneración por el trabajo que realiza, según los datos de la Encovi 2014. \n\n'\
        +' El porcentaje para el departamento es ' + self.mayor_menor(datos[2][1], datos[1][1])\
        +' que el dato nacional al ubicarse en ' + self.formato_bonito(datos[2][1])+'\\%.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','5_06.tex'), 'w')
        archivo.write(des)


    def des_507(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','5_08.csv'))
        des = 'En el 2014 el ' + self.formato_bonito(datos[3][1]) + '\\% de los '\
        +' ocupados del departamento de ' + self.lugar_geografico + ' estaban '\
        +' afilidados al seguro social. \n\n'\
        +'Este dato es ' + self.mayor_menor(datos[3][1], datos[1][1]) + ' que el '\
        +' observado en 2006, año en el que la proporción de ocupados con acceso a seguro social '\
        +' se ubicó en ' + self.formato_bonito(datos[1][1]) + '\\%.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','5_07.tex'), 'w')
        archivo.write(des)

    def des_508(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','5_09.csv'))
        des = 'Según la Encovi 2014, el ingreso promedio mensual de los ocupados en el primer empleo a nivel nacional fue de Q 2,344.5 en ese año. \n\n'\
        +'Para el departamento de ' + self.lugar_geografico + ', este indicador se ubicó en '\
        +'Q.' + self.formato_bonito(datos[2][1]) + ', monto ' + self.mayor_menor(datos[2][1],datos[1][1])\
        +' que el promedio nacional.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','5_08.tex'), 'w')
        archivo.write(des)


    def des_509(self):
        datos = self.leer_csv(os.path.join(self.ruta_salida, 'csv','5_10.csv'))
        des = 'La tasa de desempleo abierto –que representa el porcentaje de la población económicamente activa que no está ocupada y realiza acciones para conseguir un empleo- se ubicó en 2.4\\% a nivel nacional en el 2014.\n\n'\
        +'En el departamento de ' + self.lugar_geografico + ' este indicador '\
        + ' se situó en ' + self.formato_bonito( datos[2][1]) + '\\% para ese mismo '\
        +'año.'
        archivo = open( os.path.join(self.ruta_salida, 'descripciones','5_09.tex'), 'w')
        archivo.write(des)


    def mas_menos(self,dato1, dato2):
        if float(dato1) > float(dato2):
            return 'más'
        else:
            return 'menos'


    def alto_bajo(self,dato1, dato2):
        if float(dato1) > float(dato2):
            return 'alto'
        else:
            return 'bajo'


    def plural_mayor_menor(self, dato1, dato2):
        if dato2 > dato1:
            return 'mayores'
        else:
            return 'menores'

    def especial(self, dato):
        if float(dato) > 1:
            return 'mientras que el '+self.formato_bonito(dato) + '\\% poseen casas con techo de paja, palma o de un material similar. '
        else:
           return 'mientras que casi ningún hogar posee vivienda con techo de paja, palma o de un material similar.'

    def especial1(self, dato):
        if float(dato) > 1:
            return 'mientras que el '+self.formato_bonito(dato) + '\\% de hogares ocupaba una vivienda con paredes de adobe. '
        else:
           return 'mientras que casi ningún hogar habitaba una  vivienda con  paredes de adobe.'

    def formato_bonito(self, numero):
        if float(numero) < 1:
            return "{:,}".format(round(float(numero),2))
        else:
            return "{:,}".format(round(float(numero),2)).strip('0').strip('.')

    def cambio(self, dato1, dato2):
        dato1 = float(dato1)
        dato2 = float(dato2)
        if dato1 > dato2:
            return 'incremento'
        else:
            return 'descenso'

    def porcentaje(self, dato1, dato2):
        dato1 = float(dato1)
        dato2 = float(dato2)
        salida = ''
        if dato1 > dato2:
            salida = str( round((dato1 / dato2 - 1)*100,2) )
        else:
            salida = str( round((1- dato2 / dato1) * 100,2 )  )
        return salida

    def mayor_menor(self,dato1,dato2):
        if dato1 > dato2:
            return ' mayor '
        elif(dato1 < dato2):
            return ' menor '
        else:
            return ' igual '

    def encima_debajo(self,dato1,dato2):
        if dato1 > dato2:
            return ' encima '
        elif(dato1 < dato2):
            return ' debajo '
        else:
            return ' igual '

    def escribir_descripciones(self):
        self.des_101()
        self.des_102()
        self.des_103()
        self.des_104()
        self.des_105()
        self.des_106()
        self.des_107()
        self.des_108()
        self.des_109()
        self.des_110()
        self.des_111()
        self.des_112()
        self.des_113()
        self.des_114()
        self.des_115()
        self.des_116()
        self.des_117()
        self.des_118()
        self.des_119()
        self.des_120()
        self.des_121()
        self.des_122()
        self.des_123()
        self.des_201()
        self.des_202()
        self.des_203()
        self.des_204()
        self.des_205()
        self.des_206()
        self.des_207()
        self.des_208()
        self.des_209()
        self.des_210()
        self.des_211()
        self.des_212()
        self.des_213()
        self.des_214()
        self.des_215()
        self.des_216()
        self.des_217()
        self.des_218()
        self.des_219()
        self.des_601()
        self.des_602()
        self.des_603()
        self.des_604()
        self.des_605()
        self.des_606()
        self.des_607()
        self.des_608()
        self.des_609()
        self.des_610()
        self.des_611()
        self.des_612()
        self.des_613()
        self.des_301()
        self.des_302()
        self.des_303()
        self.des_304()
        self.des_305()
        self.des_401()
        self.des_402()
        self.des_403()
        self.des_404()
        self.des_405()
        self.des_501()
        self.des_502()
        self.des_503()
        self.des_504()
        self.des_505()
        self.des_506()
        self.des_507()
        self.des_508()
        self.des_509()
